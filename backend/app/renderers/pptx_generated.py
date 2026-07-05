from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.layouts.theme import get_theme
from app.schemas.slide import DeckIR, LayoutName, SlideIR, validate_image_file


SLIDE_W = 13.333
SLIDE_H = 7.5
TITLE_PT = 27
BODY_PT = 13.5
HIERARCHY_TOP_PT = 19.5
CAPTION_PT = 13.5
FOOTNOTE_PT = 13.5
MIN_BODY_PT = BODY_PT


class GeneratedPptxRenderer:
    def __init__(self, project_root: Path, max_image_bytes: int) -> None:
        self.project_root = project_root
        self.max_image_bytes = max_image_bytes

    def render(self, deck: DeckIR, output_path: Path) -> list[str]:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        warnings: list[str] = []
        for slide_ir in deck.slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._draw_background(slide, slide_ir)
            if slide_ir.layout == LayoutName.architecture_flow:
                self._architecture_flow(slide, slide_ir)
            elif slide_ir.layout == LayoutName.text_image:
                self._text_image(slide, slide_ir, warnings)
            elif slide_ir.layout == LayoutName.table:
                self._table(slide, slide_ir)
            else:
                self._generic(slide, slide_ir)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        return warnings

    def _rgb(self, value: str) -> RGBColor:
        value = value.strip("#")
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

    def _draw_background(self, slide, slide_ir: SlideIR) -> None:
        theme = get_theme(slide_ir.theme.value)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = self._rgb(theme.background)
        self._text(slide, slide_ir.title, 0.55, 0.28, 9.7, 0.62, TITLE_PT, theme.foreground, bold=True, name="title")
        if slide_ir.subtitle:
            self._text(slide, slide_ir.subtitle, 0.58, 0.98, 10.8, 0.42, BODY_PT, theme.muted, name="subtitle")
        self._text(slide, slide_ir.slide_id, 10.0, 6.86, 2.7, 0.32, FOOTNOTE_PT, theme.muted, align=PP_ALIGN.RIGHT, name="slide_key")

    def _text(self, slide, text: str, x: float, y: float, w: float, h: float, size: int, color: str, bold: bool = False, align=None, name: str | None = None):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        if name:
            box.name = name
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        p = frame.paragraphs[0]
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = text or ""
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Aptos"
        run.font.color.rgb = self._rgb(color)
        return box

    def _body_text(self, slide, text: str, x: float, y: float, w: float, h: float, color: str, name: str = "body"):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        box.name = name
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        lines = [line.rstrip() for line in (text or "").splitlines() if line.strip()]
        if not lines:
            lines = [""]
        for index, raw_line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            stripped = raw_line.strip()
            level = self._hierarchy_level(raw_line)
            clean = self._strip_bullet_marker(stripped)
            paragraph.level = min(level, 4)
            self._apply_paragraph_layout(paragraph, level, index)
            run = paragraph.add_run()
            run.text = self._display_hierarchy_text(clean, level)
            run.font.size = Pt(self._hierarchy_font_size(level))
            run.font.bold = level == 0
            run.font.name = "Aptos"
            run.font.color.rgb = self._rgb(color)
        return box

    def _hierarchy_level(self, raw_line: str) -> int:
        stripped = raw_line.lstrip()
        indent = len(raw_line) - len(stripped)
        marker_level = 1 if stripped.startswith(("・", "-", "*")) else 0
        return min(4, max(marker_level, indent // 2))

    def _strip_bullet_marker(self, stripped: str) -> str:
        if stripped.startswith(("・", "-", "*")):
            return stripped[1:].strip()
        return stripped

    def _display_hierarchy_text(self, text: str, level: int) -> str:
        if level == 0:
            return text
        return f"- {text}"

    def _hierarchy_font_size(self, level: int) -> float:
        if level == 0:
            return HIERARCHY_TOP_PT
        if level == 1:
            return BODY_PT
        return max(MIN_BODY_PT, BODY_PT - (level - 1) * 0.75)

    def _apply_paragraph_layout(self, paragraph, level: int, index: int) -> None:
        paragraph.line_spacing = 1.08
        paragraph.space_after = Pt(2)
        if level == 0:
            paragraph.space_before = Pt(9 if index else 0)
            paragraph.left_margin = Inches(0)
            paragraph.first_line_indent = Inches(0)
            return
        paragraph.space_before = Pt(3)
        paragraph.left_margin = Inches(0.28 * level)
        paragraph.first_line_indent = Inches(-0.12)

    def _panel(self, slide, x: float, y: float, w: float, h: float, fill: str, line: str, radius=True, name: str | None = None):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        if name:
            shape.name = name
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill)
        shape.line.color.rgb = self._rgb(line)
        shape.line.width = Pt(1.2)
        return shape

    def _architecture_flow(self, slide, slide_ir: SlideIR) -> None:
        theme = get_theme(slide_ir.theme.value)
        diagram = slide_ir.diagram
        if diagram is None:
            return
        count = len(diagram.nodes)
        usable_w = 11.7
        node_w = min(2.25, (usable_w - max(count - 1, 0) * 0.45) / max(count, 1))
        start_x = 0.8
        y = 2.82
        positions: dict[str, tuple[float, float, float, float]] = {}
        for index, node in enumerate(diagram.nodes):
            x = start_x + index * (node_w + 0.45)
            positions[node.id] = (x, y, node_w, 0.9)
            self._panel(slide, x, y, node_w, 0.9, theme.panel, theme.line, name=f"diagram_node_{node.id}")
            self._text(slide, node.label, x + 0.14, y + 0.16, node_w - 0.28, 0.38, BODY_PT, theme.foreground, bold=True, align=PP_ALIGN.CENTER, name=f"diagram_label_{node.id}")
            self._text(slide, node.kind, x + 0.14, y + 0.58, node_w - 0.28, 0.22, CAPTION_PT, theme.muted, align=PP_ALIGN.CENTER, name=f"diagram_kind_{node.id}")
        for edge in diagram.edges:
            if edge.source not in positions or edge.target not in positions:
                continue
            sx, sy, sw, sh = positions[edge.source]
            tx, ty, _, th = positions[edge.target]
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(sx + sw),
                Inches(sy + sh / 2),
                Inches(tx),
                Inches(ty + th / 2),
            )
            line.line.color.rgb = self._rgb(theme.accent)
            line.line.width = Pt(2)
            line.line.end_arrowhead = True
        self._body_text(slide, slide_ir.body or "", 0.8, 4.62, 11.2, 1.9, theme.muted)

    def _text_image(self, slide, slide_ir: SlideIR, warnings: list[str]) -> None:
        theme = get_theme(slide_ir.theme.value)
        body = slide_ir.body or slide_ir.subtitle or ""
        self._panel(slide, 0.75, 1.55, 5.2, 4.75, theme.panel, theme.line, radius=False, name="body_panel")
        self._body_text(slide, body, 1.05, 1.85, 4.55, 3.95, theme.foreground)
        self._panel(slide, 6.4, 1.55, 5.95, 4.75, "EEF5FF", theme.line, radius=False, name="image_panel")
        if slide_ir.image:
            try:
                image_path = validate_image_file(slide_ir.image, self.project_root, self.max_image_bytes)
                slide.shapes.add_picture(str(image_path), Inches(6.55), Inches(1.7), width=Inches(5.65), height=Inches(4.45))
            except ValueError as exc:
                warnings.append(str(exc))
                self._text(slide, "Image unavailable", 7.6, 3.5, 3.2, 0.4, BODY_PT, theme.muted, align=PP_ALIGN.CENTER)
        else:
            self._text(slide, "Image placeholder", 7.6, 3.5, 3.2, 0.4, BODY_PT, theme.muted, align=PP_ALIGN.CENTER)

    def _table(self, slide, slide_ir: SlideIR) -> None:
        theme = get_theme(slide_ir.theme.value)
        if not slide_ir.table:
            return
        data = [slide_ir.table.headers] + slide_ir.table.rows
        rows = len(data)
        cols = len(slide_ir.table.headers)
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.75), Inches(1.65), Inches(11.85), Inches(min(4.95, 0.45 * rows + 0.35)))
        table_shape.name = "table_area"
        table = table_shape.table
        for row_index, row in enumerate(data):
            for col_index, value in enumerate(row):
                cell = table.cell(row_index, col_index)
                cell.text = str(value)
                cell.fill.solid()
                if row_index == 0:
                    cell.fill.fore_color.rgb = self._rgb("1F4E79")
                elif row_index % 2 == 1:
                    cell.fill.fore_color.rgb = self._rgb("F2F2F2")
                else:
                    cell.fill.fore_color.rgb = self._rgb("FFFFFF")
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(BODY_PT)
                    paragraph.font.bold = row_index == 0
                    paragraph.font.color.rgb = self._rgb("FFFFFF" if row_index == 0 else "202A3A")
                    for run in paragraph.runs:
                        run.font.size = Pt(BODY_PT)
                        run.font.bold = row_index == 0
                        run.font.name = "Aptos"
                        run.font.color.rgb = self._rgb("FFFFFF" if row_index == 0 else "202A3A")

    def _generic(self, slide, slide_ir: SlideIR) -> None:
        theme = get_theme(slide_ir.theme.value)
        self._panel(slide, 0.85, 1.48, 11.6, 5.08, theme.panel, theme.line, radius=False, name="body_panel")
        self._body_text(slide, slide_ir.body or slide_ir.subtitle or "", 1.2, 1.76, 10.85, 4.54, theme.foreground)
