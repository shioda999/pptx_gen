from __future__ import annotations

from pathlib import Path

from pptx import Presentation


EMU_PER_INCH = 914400


def validate_pptx_layout(path: Path) -> tuple[list[str], list[str]]:
    errors: list[str] = []
    warnings: list[str] = []
    if not path.exists():
        return [f"file does not exist: {path}"], warnings
    prs = Presentation(path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    for slide_index, slide in enumerate(prs.slides, start=1):
        if len(slide.shapes) == 0:
            warnings.append(f"slide {slide_index} is empty")
        for shape in slide.shapes:
            x, y, w, h = shape.left, shape.top, shape.width, shape.height
            name = getattr(shape, "name", "shape")
            if x < 0 or y < 0 or x + w > slide_w or y + h > slide_h:
                warnings.append(f"slide {slide_index}: {name} may be outside slide bounds")
            if getattr(shape, "has_text_frame", False) and (w < 0.25 * EMU_PER_INCH or h < 0.12 * EMU_PER_INCH):
                warnings.append(f"slide {slide_index}: {name} text box is very small")
        boxes = [
            (shape.left, shape.top, shape.width, shape.height, getattr(shape, "name", "shape"))
            for shape in slide.shapes
            if shape.width > 0
            and shape.height > 0
            and not _is_slide_background(shape, slide_w, slide_h)
        ]
        for i, first in enumerate(boxes):
            for second in boxes[i + 1 :]:
                if _contains(first, second) or _contains(second, first):
                    continue
                overlap = _overlap_area(first, second)
                first_area = first[2] * first[3]
                second_area = second[2] * second[3]
                smaller = min(first_area, second_area)
                if smaller and overlap / smaller > 0.85:
                    warnings.append(f"slide {slide_index}: {first[4]} heavily overlaps {second[4]}")
    return errors, warnings


def _overlap_area(a, b) -> int:
    ax1, ay1, aw, ah, _ = a
    bx1, by1, bw, bh, _ = b
    ax2, ay2 = ax1 + aw, ay1 + ah
    bx2, by2 = bx1 + bw, by1 + bh
    width = max(0, min(ax2, bx2) - max(ax1, bx1))
    height = max(0, min(ay2, by2) - max(ay1, by1))
    return width * height


def _contains(a, b) -> bool:
    ax1, ay1, aw, ah, _ = a
    bx1, by1, bw, bh, _ = b
    return ax1 <= bx1 and ay1 <= by1 and ax1 + aw >= bx1 + bw and ay1 + ah >= by1 + bh


def _is_slide_background(shape, slide_w: int, slide_h: int) -> bool:
    return shape.left == 0 and shape.top == 0 and shape.width >= slide_w and shape.height >= slide_h
