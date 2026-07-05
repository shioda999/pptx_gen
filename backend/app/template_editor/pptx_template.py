from __future__ import annotations

import json
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from app.core.config import Settings, get_settings
from app.core.paths import ensure_output_path, resolve_project_path
from app.schemas.slide import DeckIR, SlideIR, TableSpec, TemplateSlideType, TemplateSlideVariant, validate_image_file
from app.services.speaker_notes import write_speaker_notes


TITLE_PT = 27
BODY_PT = 13.5
HIERARCHY_TOP_PT = 19.5
CAPTION_PT = 13.5
MIN_BODY_PT = BODY_PT


@dataclass
class TemplateEditResult:
    pptx_path: Path
    speaker_notes_path: Path | None
    slide_count: int
    warnings: list[str]


@dataclass(frozen=True)
class TableStyle:
    header_fill: RGBColor = RGBColor(31, 78, 121)
    odd_fill: RGBColor = RGBColor(242, 242, 242)
    even_fill: RGBColor = RGBColor(255, 255, 255)
    header_font: RGBColor = RGBColor(255, 255, 255)
    body_font: RGBColor = RGBColor(32, 42, 58)


class TemplatePptxEditor:
    """Conservative python-pptx template editor.

    This preserves the source deck and only updates shapes addressed by
    PowerPoint Selection Pane names such as title, subtitle, body, table_area,
    hero_image, and slide_key. COM automation can replace this class later.
    """

    def __init__(self, settings: Settings | None = None) -> None:
        self.settings = settings or get_settings()

    def inspect(self, path_value: str) -> dict[str, Any]:
        path = resolve_project_path(path_value, self.settings)
        prs = Presentation(path)
        metadata = self._load_metadata(path, prs)
        category_by_index = {slide["index"]: slide["category"] for slide in metadata["slides"]}
        variant_by_index = {slide["index"]: slide.get("variant", "generic") for slide in metadata["slides"]}
        slides: list[dict[str, Any]] = []
        for index, slide in enumerate(prs.slides):
            names = [shape.name for shape in slide.shapes if shape.name]
            text_shapes = [
                {
                    "name": shape.name,
                    "text": shape.text_frame.text.strip(),
                }
                for shape in slide.shapes
                if shape.name and getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
            ]
            table_candidates = [shape.name for shape in slide.shapes if getattr(shape, "has_table", False)]
            image_candidates = [
                shape.name
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or "image" in shape.name.lower()
            ]
            slides.append(
                {
                    "index": index,
                    "category": category_by_index.get(index, "content"),
                    "variant": variant_by_index.get(index, "generic"),
                    "slide_key": self._slide_key(slide),
                    "layout": slide.slide_layout.name if slide.slide_layout else None,
                    "editable_shapes": names,
                    "editable_text_shapes": text_shapes,
                    "table_candidates": table_candidates,
                    "image_candidates": image_candidates,
                }
            )
        return {"pptx_path": str(path), "metadata_path": str(self._metadata_path(path)), "slide_count": len(prs.slides), "slides": slides}

    def apply_deck(self, template_path: str, deck: DeckIR, output_filename: str) -> TemplateEditResult:
        output_path = self._copy_to_output(template_path, output_filename)
        prs = Presentation(output_path)
        warnings: list[str] = []
        source_path = resolve_project_path(template_path, self.settings)
        selected_indices, select_warnings = self._select_slide_indices(source_path, prs, deck)
        warnings.extend(select_warnings)
        self._keep_slides(prs, selected_indices)
        for index, slide_ir in enumerate(deck.slides):
            slide = prs.slides[index] if index < len(prs.slides) else None
            if slide is None:
                warnings.append(f"slide not found for {slide_ir.slide_type.value}:{slide_ir.slide_id}")
                continue
            warnings.extend(self._apply_slide_ir(slide, slide_ir))
        prs.save(output_path)
        notes_path = output_path.with_suffix(".speaker_notes.md")
        write_speaker_notes(deck, notes_path)
        return TemplateEditResult(output_path, notes_path, len(prs.slides), warnings)

    def _metadata_path(self, template_path: Path) -> Path:
        return template_path.with_suffix(".metadata.json")

    def _load_metadata(self, template_path: Path, prs: Presentation) -> dict[str, Any]:
        metadata_path = self._metadata_path(template_path)
        if metadata_path.exists():
            return json.loads(metadata_path.read_text(encoding="utf-8"))
        return self._infer_metadata(template_path, prs)

    def _infer_metadata(self, template_path: Path, prs: Presentation) -> dict[str, Any]:
        slides: list[dict[str, Any]] = []
        for index, slide in enumerate(prs.slides):
            layout_name = slide.slide_layout.name if slide.slide_layout else ""
            lowered = layout_name.lower()
            if index == 0 or "タイトルのみ" in layout_name or "title slide" in lowered:
                category = "title"
            elif "目次" in layout_name or "agenda" in lowered:
                category = "agenda"
            elif "セクション" in layout_name or "サブタイトル" in layout_name or "section" in lowered:
                category = "section"
            else:
                category = "content"
            variant = self._infer_variant(layout_name)
            slides.append({"index": index, "category": category, "variant": variant, "layout": layout_name, "notes": "auto-inferred"})
        return {"template": template_path.name, "version": 1, "slides": slides}

    def _infer_variant(self, layout_name: str) -> str:
        lowered = layout_name.lower()
        if "表" in layout_name or "table" in lowered:
            return TemplateSlideVariant.table.value
        if "2 段" in layout_name or "two" in lowered:
            return TemplateSlideVariant.two_column.value
        if "画像" in layout_name or "picture" in lowered or "image" in lowered:
            return TemplateSlideVariant.image.value
        if "タイトル" in layout_name or "title" in lowered:
            return TemplateSlideVariant.cover.value
        return TemplateSlideVariant.text.value

    def _select_slide_indices(self, template_path: Path, prs: Presentation, deck: DeckIR) -> tuple[list[int], list[str]]:
        metadata = self._load_metadata(template_path, prs)
        by_category: dict[str, list[int]] = {item.value: [] for item in TemplateSlideType}
        by_category_variant: dict[tuple[str, str], list[int]] = {}
        for slide in metadata.get("slides", []):
            category = slide.get("category", "content")
            if category not in by_category:
                category = "content"
            variant = slide.get("variant", "generic")
            if variant not in {item.value for item in TemplateSlideVariant}:
                variant = "generic"
            index = int(slide["index"])
            if 0 <= index < len(prs.slides):
                by_category[category].append(index)
                by_category_variant.setdefault((category, variant), []).append(index)
        used: set[int] = set()
        selected: list[int] = []
        warnings: list[str] = []
        for slide_ir in deck.slides:
            requested = slide_ir.slide_type.value
            requested_variant = slide_ir.slide_variant.value
            index = self._take_first_unused(by_category_variant.get((requested, requested_variant), []), used)
            if index is None and requested_variant != TemplateSlideVariant.generic.value:
                index = self._take_first_unused(by_category_variant.get((requested, TemplateSlideVariant.generic.value), []), used)
            if index is None:
                index = self._take_first_unused(by_category.get(requested, []), used)
            if index is None and requested != TemplateSlideType.content.value:
                index = self._take_first_unused(by_category_variant.get((TemplateSlideType.content.value, requested_variant), []), used)
                if index is None:
                    index = self._take_first_unused(by_category.get(TemplateSlideType.content.value, []), used)
                if index is not None:
                    warnings.append(f"no unused {requested}/{requested_variant} template slide; used content slide for {slide_ir.slide_id}")
            if index is None:
                index = self._take_first_unused(list(range(len(prs.slides))), used)
                if index is not None:
                    warnings.append(f"no unused template slide for {slide_ir.slide_id}; used slide {index}")
            if index is None:
                warnings.append(f"not enough template slides for {slide_ir.slide_id}")
                continue
            used.add(index)
            selected.append(index)
        return selected, warnings

    def _take_first_unused(self, candidates: list[int], used: set[int]) -> int | None:
        for index in candidates:
            if index not in used:
                return index
        return None

    def _keep_slides(self, prs: Presentation, selected_indices: list[int]) -> None:
        slide_id_list = prs.slides._sldIdLst
        original = list(slide_id_list)
        selected = [original[index] for index in selected_indices]
        for element in original:
            slide_id_list.remove(element)
        for element in selected:
            slide_id_list.append(element)

    def update_shapes(self, pptx_path: str, output_filename: str | None, slide_key: str | None, slide_index: int | None, values: dict[str, Any]) -> TemplateEditResult:
        output_path = self._copy_to_output(pptx_path, output_filename or Path(pptx_path).name)
        prs = Presentation(output_path)
        slide = self._find_slide(prs, slide_key, fallback_index=slide_index)
        if slide is None:
            raise ValueError("target slide not found")
        warnings: list[str] = []
        for shape_name, value in values.items():
            shape = self._find_shape(slide, shape_name)
            if shape is None:
                warnings.append(f"shape not found: {shape_name}")
                continue
            if getattr(shape, "has_text_frame", False):
                role = self._role_from_shape_name(shape_name)
                if role:
                    self._set_text_frame(shape, str(value), role)
                else:
                    shape.text_frame.text = str(value)
            else:
                warnings.append(f"shape is not text-editable: {shape_name}")
        prs.save(output_path)
        return TemplateEditResult(output_path, None, len(prs.slides), warnings)

    def _role_from_shape_name(self, shape_name: str) -> str | None:
        lowered = shape_name.lower()
        if lowered == "title":
            return "title"
        if lowered == "subtitle":
            return "body"
        if lowered == "body":
            return "hierarchy"
        return None

    def set_table(self, pptx_path: str, output_filename: str | None, slide_key: str | None, slide_index: int | None, shape_name: str, table: TableSpec) -> TemplateEditResult:
        output_path = self._copy_to_output(pptx_path, output_filename or Path(pptx_path).name)
        prs = Presentation(output_path)
        slide = self._find_slide(prs, slide_key, fallback_index=slide_index)
        if slide is None:
            raise ValueError("target slide not found")
        warnings = self._replace_table(slide, shape_name, table)
        prs.save(output_path)
        return TemplateEditResult(output_path, None, len(prs.slides), warnings)

    def set_image(self, pptx_path: str, output_filename: str | None, slide_key: str | None, slide_index: int | None, shape_name: str, image_path: str) -> TemplateEditResult:
        output_path = self._copy_to_output(pptx_path, output_filename or Path(pptx_path).name)
        prs = Presentation(output_path)
        slide = self._find_slide(prs, slide_key, fallback_index=slide_index)
        if slide is None:
            raise ValueError("target slide not found")
        warnings = self._replace_image(slide, shape_name, image_path)
        prs.save(output_path)
        return TemplateEditResult(output_path, None, len(prs.slides), warnings)

    def _copy_to_output(self, source_value: str, output_filename: str) -> Path:
        source = resolve_project_path(source_value, self.settings)
        if not source.exists():
            raise ValueError(f"template/pptx file does not exist: {source_value}")
        if source.suffix.lower() != ".pptx":
            raise ValueError("template/pptx file must be .pptx")
        output_path = ensure_output_path(output_filename, self.settings)
        if output_path.suffix.lower() != ".pptx":
            output_path = output_path.with_suffix(".pptx")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        if source.resolve() != output_path.resolve():
            shutil.copyfile(source, output_path)
        return output_path

    def _apply_slide_ir(self, slide, slide_ir: SlideIR) -> list[str]:
        warnings: list[str] = []
        used_shape_ids: set[int] = set()
        title_shape = self._find_shape(slide, "title") or self._find_text_shape(slide, ["タイトル", "Title"])
        subtitle_shape = self._find_shape(slide, "subtitle") or self._find_text_shape(slide, ["サブタイトル", "Subtitle"])
        body_shape = self._find_shape(slide, "body") or self._find_text_shape(
            slide,
            ["body", "本文", "コンテンツ プレースホルダー", "テキスト プレースホルダー", "Content Placeholder", "Text Placeholder"],
            exclude=[title_shape, subtitle_shape],
        )
        for shape, value, role in (
            (title_shape, slide_ir.title, "title"),
            (subtitle_shape, slide_ir.subtitle or "", "body"),
            (body_shape, slide_ir.body or "", "hierarchy"),
        ):
            if shape is not None and getattr(shape, "has_text_frame", False) and value:
                self._set_text_frame(shape, value, role)
                used_shape_ids.add(id(shape))
        if slide_ir.table:
            warnings.extend(self._replace_table(slide, "table_area", slide_ir.table))
        if slide_ir.diagram:
            warnings.extend(self._replace_diagram(slide, "diagram_area", slide_ir))
        if slide_ir.image:
            warnings.extend(self._replace_image(slide, slide_ir.image.target, slide_ir.image.path))
        self._clear_known_sample_text(slide, used_shape_ids)
        return warnings

    def _replace_diagram(self, slide, shape_name: str, slide_ir: SlideIR) -> list[str]:
        target = self._find_shape(slide, shape_name) or self._find_content_area(slide)
        if target is None or slide_ir.diagram is None:
            return [f"diagram target not found: {shape_name}"]
        left, top, width, height = target.left, target.top, target.width, target.height
        self._remove_shape(target)
        nodes = slide_ir.diagram.nodes
        count = len(nodes)
        gap = int(width * 0.035)
        node_w = int((width - gap * max(0, count - 1)) / max(1, count))
        node_h = int(height * 0.34)
        node_top = int(top + height * 0.31)
        positions: dict[str, tuple[int, int, int, int]] = {}
        for index, node in enumerate(nodes):
            x = int(left + index * (node_w + gap))
            positions[node.id] = (x, node_top, node_w, node_h)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, node_top, node_w, node_h)
            box.name = f"diagram_node_{node.id}"
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            box.line.color.rgb = RGBColor(45, 140, 255)
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = node.label
            r.font.size = Pt(BODY_PT)
            r.font.bold = True
            r.font.color.rgb = RGBColor(24, 32, 51)
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = node.kind
            r2.font.size = Pt(CAPTION_PT)
            r2.font.color.rgb = RGBColor(82, 96, 113)
        for edge in slide_ir.diagram.edges:
            if edge.source not in positions or edge.target not in positions:
                continue
            sx, sy, sw, sh = positions[edge.source]
            tx, ty, _, th = positions[edge.target]
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx + sw, sy + sh // 2, tx, ty + th // 2)
            line.name = f"diagram_edge_{edge.source}_{edge.target}"
            line.line.color.rgb = RGBColor(45, 140, 255)
            line.line.width = Pt(2)
            line.line.end_arrowhead = True
        return []

    def _replace_table(self, slide, shape_name: str, table: TableSpec) -> list[str]:
        target = self._find_shape(slide, shape_name) or self._find_table_area(slide) or self._find_content_area(slide)
        if target is None:
            return [f"table target not found: {shape_name}"]
        if getattr(target, "has_table", False):
            style = self._extract_table_style(target)
            return self._fill_existing_table(target, table, style)
        left, top, width, height = target.left, target.top, target.width, target.height
        self._remove_shape(target)
        rows = [table.headers] + table.rows
        table_shape = slide.shapes.add_table(len(rows), len(table.headers), left, top, width, height)
        table_shape.name = shape_name
        self._style_generated_table(table_shape, TableStyle())
        ppt_table = table_shape.table
        for row_index, row in enumerate(rows):
            for col_index, value in enumerate(row):
                cell = ppt_table.cell(row_index, col_index)
                cell.text = str(value)
                self._format_table_cell(cell, row_index, TableStyle())
        return []

    def _fill_existing_table(self, table_shape, table: TableSpec, style: TableStyle) -> list[str]:
        ppt_table = table_shape.table
        data = [table.headers] + table.rows
        warnings: list[str] = []
        self._shrink_table_to_fit(ppt_table, len(data), len(table.headers))
        if len(data) > len(ppt_table.rows):
            warnings.append(f"template table has {len(ppt_table.rows)} rows; truncated {len(data) - len(ppt_table.rows)} rows")
        if len(table.headers) > len(ppt_table.columns):
            warnings.append(f"template table has {len(ppt_table.columns)} columns; truncated {len(table.headers) - len(ppt_table.columns)} columns")
        for row_index in range(len(ppt_table.rows)):
            for col_index in range(len(ppt_table.columns)):
                cell = ppt_table.cell(row_index, col_index)
                if row_index < len(data) and col_index < len(data[row_index]):
                    cell.text = str(data[row_index][col_index])
                else:
                    cell.text = ""
                self._apply_banded_table_cell_style(cell, row_index, style)
        return warnings

    def _shrink_table_to_fit(self, ppt_table, desired_rows: int, desired_cols: int) -> None:
        table_xml = ppt_table._tbl
        while len(table_xml.tr_lst) > desired_rows:
            row = table_xml.tr_lst[-1]
            row.getparent().remove(row)
        while len(table_xml.tblGrid.gridCol_lst) > desired_cols:
            col = table_xml.tblGrid.gridCol_lst[-1]
            col.getparent().remove(col)
            for row in table_xml.tr_lst:
                if len(row.tc_lst) > desired_cols:
                    cell = row.tc_lst[-1]
                    cell.getparent().remove(cell)

    def _style_generated_table(self, table_shape, style: TableStyle) -> None:
        table = table_shape.table
        for row_index, row in enumerate(table.rows):
            for col_index, cell in enumerate(row.cells):
                self._apply_banded_table_cell_style(cell, row_index, style)

    def _apply_banded_table_cell_style(self, cell, row_index: int, style: TableStyle) -> None:
        cell.fill.solid()
        if row_index == 0:
            cell.fill.fore_color.rgb = style.header_fill
        elif row_index % 2 == 1:
            cell.fill.fore_color.rgb = style.odd_fill
        else:
            cell.fill.fore_color.rgb = style.even_fill
        self._format_table_cell(cell, row_index, style)

    def _format_table_cell(self, cell, row_index: int, style: TableStyle) -> None:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(BODY_PT)
            paragraph.font.bold = row_index == 0
            paragraph.font.name = "Aptos"
            if row_index == 0:
                paragraph.font.color.rgb = style.header_font
            else:
                paragraph.font.color.rgb = style.body_font
            for run in paragraph.runs:
                run.font.size = Pt(BODY_PT)
                run.font.bold = row_index == 0
                run.font.name = "Aptos"
                if row_index == 0:
                    run.font.color.rgb = style.header_font
                else:
                    run.font.color.rgb = style.body_font

    def _extract_table_style(self, table_shape) -> TableStyle:
        table = table_shape.table
        default = TableStyle()
        header_fill = self._cell_fill_rgb(table.cell(0, 0)) or default.header_fill
        odd_fill = self._cell_fill_rgb(table.cell(1, 0)) if len(table.rows) > 1 else None
        even_fill = self._cell_fill_rgb(table.cell(2, 0)) if len(table.rows) > 2 else None
        header_font = self._cell_font_rgb(table.cell(0, 0)) or self._contrasting_font(header_fill)
        body_font = self._cell_font_rgb(table.cell(1, 0)) if len(table.rows) > 1 else None
        return TableStyle(
            header_fill=header_fill,
            odd_fill=odd_fill or default.odd_fill,
            even_fill=even_fill or default.even_fill,
            header_font=header_font,
            body_font=body_font or default.body_font,
        )

    def _cell_fill_rgb(self, cell) -> RGBColor | None:
        try:
            if cell.fill.type is not None and cell.fill.fore_color.type is not None:
                return cell.fill.fore_color.rgb
        except (AttributeError, TypeError):
            return None
        return None

    def _cell_font_rgb(self, cell) -> RGBColor | None:
        try:
            for paragraph in cell.text_frame.paragraphs:
                if paragraph.font.color.type is not None:
                    return paragraph.font.color.rgb
        except (AttributeError, TypeError):
            return None
        return None

    def _contrasting_font(self, fill: RGBColor) -> RGBColor:
        brightness = (fill[0] * 299 + fill[1] * 587 + fill[2] * 114) / 1000
        return RGBColor(255, 255, 255) if brightness < 140 else RGBColor(32, 42, 58)

    def _set_text_frame(self, shape, text: str, role: str) -> None:
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        if role == "hierarchy":
            lines = [line.rstrip() for line in text.splitlines() if line.strip()]
        else:
            lines = [text]
        if not lines:
            lines = [""]
        for index, raw_line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            stripped = raw_line.strip()
            level = self._hierarchy_level(raw_line) if role == "hierarchy" else 0
            run = paragraph.add_run()
            clean = self._strip_bullet_marker(stripped) if role == "hierarchy" else stripped
            run.text = self._display_hierarchy_text(clean, level) if role == "hierarchy" else clean
            run.font.name = "Aptos"
            if role == "title":
                run.font.size = Pt(TITLE_PT)
                run.font.bold = True
            elif role == "hierarchy":
                run.font.size = Pt(self._hierarchy_font_size(level))
                run.font.bold = level == 0
                paragraph.level = min(level, 4)
                self._apply_paragraph_layout(paragraph, level, index)
            else:
                run.font.size = Pt(BODY_PT)
                run.font.bold = False

    def _hierarchy_level(self, raw_line: str) -> int:
        stripped = raw_line.lstrip()
        indent = len(raw_line) - len(stripped)
        marker_level = 1 if stripped.startswith(("・", "-", "*")) else 0
        return min(4, max(marker_level, indent // 2))

    def _strip_bullet_marker(self, stripped: str) -> str:
        if stripped.startswith(("・", "-", "*")):
            return stripped[1:].strip()
        return stripped

    def _display_hierarchy_text(self, text: str, level: int) -> str:
        if level == 0:
            return text
        return f"- {text}"

    def _hierarchy_font_size(self, level: int) -> float:
        if level == 0:
            return HIERARCHY_TOP_PT
        if level == 1:
            return BODY_PT
        return max(MIN_BODY_PT, BODY_PT - (level - 1) * 0.75)

    def _apply_paragraph_layout(self, paragraph, level: int, index: int) -> None:
        paragraph.line_spacing = 1.08
        paragraph.space_after = Pt(2)
        if level == 0:
            paragraph.space_before = Pt(9 if index else 0)
            paragraph.left_margin = Pt(0)
            paragraph.first_line_indent = Pt(0)
            return
        paragraph.space_before = Pt(3)
        paragraph.left_margin = Pt(20 * level)
        paragraph.first_line_indent = Pt(-8)

    def _replace_image(self, slide, shape_name: str, image_path_value: str) -> list[str]:
        target = self._find_shape(slide, shape_name)
        if target is None:
            return [f"image target not found: {shape_name}"]
        image_path = validate_image_file(
            type("ImageProxy", (), {"path": image_path_value})(),
            self.settings.project_root,
            self.settings.max_image_bytes,
        )
        left, top, width, height = target.left, target.top, target.width, target.height
        self._remove_shape(target)
        picture = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        picture.name = shape_name
        return []

    def _find_slide(self, prs: Presentation, slide_key: str | None, fallback_index: int | None = None):
        if slide_key:
            for slide in prs.slides:
                if self._slide_key(slide) == slide_key:
                    return slide
        if fallback_index is not None and 0 <= fallback_index < len(prs.slides):
            return prs.slides[fallback_index]
        return None

    def _slide_key(self, slide) -> str | None:
        shape = self._find_shape(slide, "slide_key")
        if shape is not None and getattr(shape, "has_text_frame", False):
            return shape.text_frame.text.strip() or None
        return None

    def _find_shape(self, slide, name: str | None):
        if not name:
            return None
        for shape in slide.shapes:
            if shape.name == name:
                return shape
        return None

    def _find_text_shape(self, slide, name_fragments: list[str], exclude: list | None = None):
        exclude_ids = {id(shape) for shape in (exclude or []) if shape is not None}
        for shape in slide.shapes:
            if id(shape) in exclude_ids or not getattr(shape, "has_text_frame", False):
                continue
            if any(fragment.lower() in shape.name.lower() for fragment in name_fragments):
                return shape
        return None

    def _find_content_area(self, slide):
        candidates = [
            "table_area",
            "diagram_area",
            "コンテンツ プレースホルダー",
            "Content Placeholder",
            "表プレースホルダー",
            "Table Placeholder",
            "図プレースホルダー",
            "Picture Placeholder",
        ]
        for fragment in candidates:
            for shape in slide.shapes:
                if fragment.lower() in shape.name.lower():
                    return shape
        return None

    def _find_table_area(self, slide):
        candidates = ["table_area", "表プレースホルダー", "Table Placeholder"]
        for fragment in candidates:
            for shape in slide.shapes:
                if fragment.lower() in shape.name.lower() or getattr(shape, "has_table", False):
                    return shape
        return None

    def _clear_known_sample_text(self, slide, used_shape_ids: set[int]) -> None:
        sample_markers = [
            "lorem ipsum",
            "tailwindtraders",
            "4567 main st",
            "555-0100",
            "弊社の理念",
            "弊社のチーム",
            "お問い合わせ",
            "弊社の専門家チーム",
            "公の場で話す時",
            "声の高さ",
            "高低の変化",
            "トーンの変化",
            "声量",
            "q&a セッション",
            "落ち着きを保つ",
            "冷静になる",
            "アクティブ リスニング",
            "フィードバックを求める",
            "パフォーマンスを振り返る",
            "新しいテクニック",
            "個人の目標",
            "反復し、適応する",
            "プレゼンテーションの提供",
            "長く残る印象",
            "効果的なコミュニケーション",
            "聞き手に意欲",
        ]
        for shape in slide.shapes:
            if id(shape) in used_shape_ids or not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            lowered = text.lower()
            if any(marker in lowered for marker in sample_markers):
                shape.text_frame.clear()

    def _remove_shape(self, shape) -> None:
        shape._element.getparent().remove(shape._element)
