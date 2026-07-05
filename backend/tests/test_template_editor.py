from pathlib import Path

from app.schemas.api import CreateDeckRequest
from app.services.deck_service import DeckService
from app.template_editor.pptx_template import TemplatePptxEditor
from app.markdown.parser import parse_markdown_deck
from pptx import Presentation


def test_inspect_template_and_apply_deck():
    project_root = Path(__file__).resolve().parents[2]
    template_path = project_root / "templates" / "codex-template.pptx"
    if not template_path.exists():
        return

    inspected = TemplatePptxEditor().inspect("templates/codex-template.pptx")
    assert inspected["slide_count"] == 4
    assert inspected["slides"][0]["slide_key"] == "codex-cover"

    markdown = (project_root / "examples" / "codex-self-intro.md").read_text(encoding="utf-8")
    deck = parse_markdown_deck(markdown)
    assert deck.slides[2].slide_type.value == "content"
    assert deck.slides[2].slide_variant.value == "table"
    response = DeckService().create_deck(
        CreateDeckRequest(
            markdown=markdown,
            template_path="templates/codex-template.pptx",
            output_filename="pytest-template-based.pptx",
        )
    )
    assert Path(response.pptx_path).exists()
    assert response.slide_count == 4
    assert all("not enough template slides" not in warning for warning in response.warnings)


def test_selective_template_slides_from_metadata():
    project_root = Path(__file__).resolve().parents[2]
    template_path = project_root / "templates" / "imported" / "template3.pptx"
    metadata_path = template_path.with_suffix(".metadata.json")
    if not template_path.exists() or not metadata_path.exists():
        return

    markdown = (project_root / "examples" / "codex-self-intro.md").read_text(encoding="utf-8")
    response = DeckService().create_deck(
        CreateDeckRequest(
            markdown=markdown,
            template_path="templates/imported/template3.pptx",
            output_filename="pytest-selective-template3.pptx",
        )
    )
    assert Path(response.pptx_path).exists()
    assert response.slide_count == 4
    assert response.warnings == []


def test_generated_deck_can_be_inspected_and_edited_by_semantic_shape_names():
    project_root = Path(__file__).resolve().parents[2]
    markdown = """---
slide_id: edit-demo
mode: generated
layout: summary
theme: modern-tech
title: Original Title
body: |
  Original section
    - Original detail
---
"""
    created = DeckService().create_deck(
        CreateDeckRequest(markdown=markdown, output_filename="pytest-edit-demo.pptx")
    )

    inspected = TemplatePptxEditor().inspect(created.pptx_path)
    text_shape_names = {shape["name"] for shape in inspected["slides"][0]["editable_text_shapes"]}
    assert {"title", "body", "slide_key"}.issubset(text_shape_names)

    result = TemplatePptxEditor().update_shapes(
        created.pptx_path,
        "pytest-edit-demo-updated.pptx",
        None,
        0,
        {
            "title": "Edited Title",
            "body": "Edited section\n  - Edited detail",
        },
    )

    prs = Presentation(result.pptx_path)
    text_by_name = {
        shape.name: shape.text_frame.text
        for shape in prs.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
    }
    assert text_by_name["title"] == "Edited Title"
    assert "Edited section" in text_by_name["body"]
    assert "Edited detail" in text_by_name["body"]
    assert result.warnings == []
