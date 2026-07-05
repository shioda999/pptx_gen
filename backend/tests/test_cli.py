import json
from pathlib import Path

from pptx import Presentation

from app.cli import main


def test_cli_create_inspect_and_update(capsys):
    markdown = Path(__file__).resolve().parents[2] / "examples" / "table.md"

    assert main(["create", "--markdown-file", str(markdown), "--output", "pytest-cli-table.pptx"]) == 0
    created = json.loads(capsys.readouterr().out)
    assert created["slide_count"] == 1
    assert Path(created["pptx_path"]).exists()

    assert main(["inspect", "--pptx-path", created["pptx_path"]]) == 0
    inspected = json.loads(capsys.readouterr().out)
    shape_names = {shape["name"] for shape in inspected["slides"][0]["editable_text_shapes"]}
    assert "title" in shape_names

    assert (
        main(
            [
                "update",
                "--pptx-path",
                created["pptx_path"],
                "--output",
                "pytest-cli-table-updated.pptx",
                "--slide-index",
                "0",
                "--set",
                "title=CLI Edited Title",
            ]
        )
        == 0
    )
    updated = json.loads(capsys.readouterr().out)
    prs = Presentation(updated["pptx_path"])
    text_by_name = {
        shape.name: shape.text_frame.text
        for shape in prs.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
    }
    assert text_by_name["title"] == "CLI Edited Title"


def test_cli_validate_reports_json(capsys):
    markdown = Path(__file__).resolve().parents[2] / "examples" / "architecture-flow.md"

    assert main(["validate", "--markdown-file", str(markdown)]) == 0
    result = json.loads(capsys.readouterr().out)

    assert result["ok"] is True
    assert result["errors"] == []
